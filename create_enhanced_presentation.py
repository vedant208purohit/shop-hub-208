from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_enhanced_shop_hub_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content
    section_header_layout = prs.slide_layouts[2]  # Section header
    blank_layout = prs.slide_layouts[6]  # Blank slide for custom layouts
    
    # Define custom colors
    primary_color = RGBColor(59, 130, 246)  # Blue
    secondary_color = RGBColor(16, 185, 129)  # Green
    accent_color = RGBColor(245, 158, 11)  # Orange
    dark_color = RGBColor(31, 41, 55)  # Dark gray
    light_color = RGBColor(249, 250, 251)  # Light gray
    
    # Slide 1: Enhanced Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Shop Hub"
    subtitle.text = "Modern E-Commerce Platform\nInternship Project Report"
    
    # Apply styling to title
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Apply styling to subtitle
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = dark_color
    subtitle_frame.paragraphs[1].font.size = Pt(18)
    subtitle_frame.paragraphs[1].font.color.rgb = secondary_color
    
    # Add additional info
    p = subtitle_frame.add_paragraph()
    p.text = "Technology Stack: React, TypeScript, Vite, Tailwind CSS, Shadcn/ui"
    p.font.size = Pt(14)
    p.font.color.rgb = dark_color
    
    p = subtitle_frame.add_paragraph()
    p.text = "Project Duration: 15th May to 10th July"
    p.font.size = Pt(14)
    p.font.color.rgb = dark_color
    
    p = subtitle_frame.add_paragraph()
    p.text = "Supervisor: Mr Yash | Date: 10th July, 2025"
    p.font.size = Pt(12)
    p.font.color.rgb = dark_color
    
    # Slide 2: Project Overview with enhanced styling
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content_text = """🏪 Modern E-Commerce Web Application

• Built using React 18 and TypeScript
• Comprehensive online shopping experience
• Features: user authentication, product browsing, shopping cart, secure checkout

🎯 Key Objectives:
• Develop user-friendly e-commerce platform with modern UI/UX
• Implement secure user authentication and session management
• Create intuitive product browsing and search experience
• Build robust shopping cart system with quantity management
• Ensure responsive design for all device types"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 3: Key Features with icons
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Features"
    content_text = """🔐 User Authentication
• Secure login/signup with validation
• Session management across browser sessions

🛍️ Product Catalog
• Browse products by categories
• Search functionality
• Detailed product information with image galleries

🛒 Shopping Cart
• Add, remove, and update product quantities
• Real-time cart updates
• Persistent cart data

📱 Responsive Design
• Mobile-first approach
• Modern UI with consistent design language
• Real-time validation with immediate feedback"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 4: Technology Stack with visual elements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Stack"
    content_text = """⚛️ Frontend Framework: React 18 with TypeScript
⚡ Build Tool: Vite for fast development and building
🎨 Styling: Tailwind CSS with Shadcn/ui components
🔄 State Management: React Context API
🛣️ Routing: React Router DOM
🎯 Icons: Lucide React
✅ Validation: Custom validation utilities

💡 Benefits:
• Type safety with TypeScript
• Fast development with Vite
• Modern UI components with Shadcn/ui
• Efficient state management"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 5: Application Screenshots - Login Page
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Application Screenshots - Login Page"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add placeholder for login screenshot
    img_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    img_frame = img_box.text_frame
    img_frame.text = "[Login Page Screenshot]\n\nSecure user authentication with\nemail validation and password strength indicator"
    img_frame.paragraphs[0].font.size = Pt(18)
    img_frame.paragraphs[0].font.color.rgb = dark_color
    img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    img_frame.paragraphs[1].font.size = Pt(14)
    img_frame.paragraphs[1].font.color.rgb = secondary_color
    img_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Slide 6: Application Screenshots - Homepage
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Application Screenshots - Homepage"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add placeholder for homepage screenshot
    img_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    img_frame = img_box.text_frame
    img_frame.text = "[Homepage Screenshot]\n\nModern responsive design with\nfeatured products and category navigation"
    img_frame.paragraphs[0].font.size = Pt(18)
    img_frame.paragraphs[0].font.color.rgb = dark_color
    img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    img_frame.paragraphs[1].font.size = Pt(14)
    img_frame.paragraphs[1].font.color.rgb = secondary_color
    img_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Slide 7: Application Screenshots - Product Page
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Application Screenshots - Product Page"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add placeholder for product page screenshot
    img_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    img_frame = img_box.text_frame
    img_frame.text = "[Product Page Screenshot]\n\nDetailed product information with\nimage gallery and add to cart functionality"
    img_frame.paragraphs[0].font.size = Pt(18)
    img_frame.paragraphs[0].font.color.rgb = dark_color
    img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    img_frame.paragraphs[1].font.size = Pt(14)
    img_frame.paragraphs[1].font.color.rgb = secondary_color
    img_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Slide 8: Application Screenshots - Cart Page
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Application Screenshots - Cart Page"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add placeholder for cart page screenshot
    img_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    img_frame = img_box.text_frame
    img_frame.text = "[Cart Page Screenshot]\n\nShopping cart with quantity management\nand persistent cart data"
    img_frame.paragraphs[0].font.size = Pt(18)
    img_frame.paragraphs[0].font.color.rgb = dark_color
    img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    img_frame.paragraphs[1].font.size = Pt(14)
    img_frame.paragraphs[1].font.color.rgb = secondary_color
    img_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Slide 9: System Architecture Diagram
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add architecture diagram placeholder
    arch_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    arch_frame = arch_box.text_frame
    arch_frame.text = "[System Architecture Diagram]\n\nComponent-Based Architecture:\n\n📱 Presentation Layer (React Components)\n⚙️ Business Logic Layer (Custom Hooks)\n🔄 State Management Layer (Context API)\n💾 Data Layer (Static Data + localStorage)"
    arch_frame.paragraphs[0].font.size = Pt(18)
    arch_frame.paragraphs[0].font.color.rgb = dark_color
    arch_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    arch_frame.paragraphs[2].font.size = Pt(14)
    arch_frame.paragraphs[2].font.color.rgb = secondary_color
    arch_frame.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    # Slide 10: Functional Requirements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Functional Requirements"
    content_text = """👤 User Management
• User registration with email, name, and password
• User login/logout functionality
• Session persistence across browser sessions
• Authentication required for cart operations

📦 Product Management
• Browse all available products
• Filter products by category
• Search products by name
• View detailed product information

🛒 Shopping Cart
• Add products to cart (authenticated users)
• View cart contents and total
• Update product quantities
• Remove products and clear entire cart
• Cart persistence across sessions"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 11: Non-Functional Requirements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Non-Functional Requirements"
    content_text = """⚡ Performance
• Application loads within 3 seconds
• Smooth animations and transitions (60fps)
• Efficient state management with minimal re-renders

🔒 Security
• Secure password requirements (8+ chars, complexity)
• Email format validation
• Input sanitization and validation
• Authentication required for sensitive operations

🎨 Usability
• Intuitive navigation and user flow
• Consistent design language throughout
• Accessible design with proper contrast
• Mobile-first responsive design

🛡️ Reliability
• Graceful error handling and recovery
• Data persistence across browser sessions
• Consistent behavior across different browsers"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 12: Development Methodology
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Development Methodology"
    content_text = """🔄 Agile Development Approach
• Iterative development cycles
• Focus on delivering working features incrementally

📋 Development Process:
1. Planning Phase
   • Requirements analysis
   • Technology selection
   • Architecture planning
   • Timeline planning

2. Design Phase
   • UI/UX design with wireframes
   • Component design
   • State management design
   • API design

3. Development Phase
   • Sprint planning
   • Component development
   • Integration
   • Continuous testing

4. Testing Phase
   • Unit testing
   • Integration testing
   • User acceptance testing
   • Performance testing"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 13: Testing Strategy
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Testing Strategy"
    content_text = """🧪 Comprehensive Testing Approach

Unit Testing:
• Component functionality testing
• Utility function testing
• Custom hook testing

Integration Testing:
• Component interaction testing
• Context integration testing
• Routing functionality testing

User Acceptance Testing:
• Complete user journey testing
• Cross-browser compatibility
• Responsive functionality testing

🔍 Test Cases:
• Authentication: Valid/invalid login, registration, session persistence
• Cart Functionality: Add, update, remove items, cart persistence
• Product Browsing: Listing, filtering, search, details view
• Validation: Email, password, name validation with real-time feedback

⚡ Performance Testing:
• Load time within 3 seconds
• Smooth interactions and animations
• Efficient memory management
• Optimized JavaScript bundle"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 14: Project Achievements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Achievements"
    content_text = """✅ Complete E-commerce Functionality
• Full product browsing, cart management, and user authentication
• Modern technology stack with React 18, TypeScript, Vite, and Tailwind CSS
• Responsive design with mobile-first approach

🎨 Technical Highlights
• Component-based architecture for modular and maintainable code
• Full TypeScript implementation for type safety
• Efficient state management using React Context API
• Custom validation system with comprehensive error handling
• Modern UI/UX with accessibility considerations

💼 Business Value
• Intuitive and engaging shopping experience
• Scalable architecture supporting future feature additions
• Clean code structure for easy maintenance
• Optimized performance for fast loading and smooth interactions
• Secure user data handling and authentication"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 15: Lessons Learned
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Lessons Learned"
    content_text = """💡 Technical Lessons

State Management:
• Proper state management is crucial for complex applications
• React Context API provides excellent state management for medium-sized applications

Component Design:
• Reusable components significantly improve development efficiency
• Component composition and prop drilling alternatives are essential

TypeScript Benefits:
• TypeScript provides significant benefits in large projects
• Type safety reduces bugs and improves development experience

Performance Optimization:
• Performance considerations should be built into the development process
• Efficient re-renders and state updates are crucial for smooth UX

🔄 Development Process Lessons

Planning and Architecture:
• Proper planning and architecture design saves significant development time
• Component structure and data flow should be planned before implementation

Testing Strategy:
• Testing should be integrated into the development process
• Different types of testing are needed for different aspects of the application

User Experience:
• User experience should be prioritized from the beginning
• Small details like loading states and error handling significantly impact UX"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 16: Future Improvements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Improvements"
    content_text = """🚀 Technical Enhancements

Backend Integration:
• Implement real backend API integration
• Add proper database for product and user data
• Implement secure payment processing

Advanced Features:
• Advanced search and filtering capabilities
• Product review and rating system
• User wishlist functionality
• Complete order processing system
• Admin interface for product management

⚡ Performance Optimizations

Code Splitting:
• Implement lazy loading for better performance
• Add proper caching strategies
• Implement content delivery network
• Convert to Progressive Web Application

📱 Additional Features:
• User Reviews and Ratings
• Wishlist Functionality
• Order Management System
• Admin Panel for Product Management
• Advanced Analytics and Reporting"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Slide 17: Project Links and Demo
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Project Links & Demo"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add project links
    links_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    links_frame = links_box.text_frame
    links_frame.text = "🔗 GitHub Repository\nhttps://github.com/vedant208purohit/shop-hub-208.git\n\n🌐 Live Demo\n[Available upon request]\n\n👤 Demo Credentials\nEmail: demo@example.com\nPassword: password\n\n📱 Features Demo\n• User Authentication: Complete login/signup flow\n• Product Browsing: Category filtering and search\n• Shopping Cart: Add, update, and remove items\n• Responsive Design: Mobile and desktop experience\n• Form Validation: Real-time validation feedback"
    links_frame.paragraphs[0].font.size = Pt(18)
    links_frame.paragraphs[0].font.color.rgb = primary_color
    links_frame.paragraphs[0].font.bold = True
    links_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    links_frame.paragraphs[1].font.size = Pt(12)
    links_frame.paragraphs[1].font.color.rgb = secondary_color
    links_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Slide 18: Conclusion
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    content_text = """🎯 Project Success

The Shop Hub e-commerce application successfully demonstrates:
• Modern web development practices
• Solid foundation for production-ready e-commerce platform
• Importance of proper planning and architecture design
• User experience considerations in web application development

📚 Key Takeaways

• Lessons learned will be valuable for future development projects
• Insights into building scalable, maintainable, and user-friendly web applications
• Understanding of modern React development patterns
• Experience with TypeScript and modern build tools

🌟 Thank You!

Questions & Discussion"""
    
    content.text = content_text
    
    # Apply styling
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].font.bold = True
    
    # Save the presentation
    prs.save('Shop_Hub_Enhanced_Presentation.pptx')
    print("Enhanced PowerPoint presentation created successfully: Shop_Hub_Enhanced_Presentation.pptx")

if __name__ == "__main__":
    create_enhanced_shop_hub_presentation() 