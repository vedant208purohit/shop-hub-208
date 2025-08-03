from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

def create_shop_hub_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content
    section_header_layout = prs.slide_layouts[2]  # Section header
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Shop Hub"
    subtitle.text = "Modern E-Commerce Platform\nInternship Project Report\n\nTechnology Stack: React, TypeScript, Vite, Tailwind CSS, Shadcn/ui\nProject Duration: 15th May to 10th July\n\nSupervisor: Mr Yash\nDate: 10th July, 2025"
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content_text = """• Modern, responsive e-commerce web application
• Built using React and TypeScript
• Comprehensive online shopping experience
• Features: user authentication, product browsing, shopping cart, secure checkout

Key Objectives:
• Develop user-friendly e-commerce platform with modern UI/UX
• Implement secure user authentication and session management
• Create intuitive product browsing and search experience
• Build robust shopping cart system with quantity management
• Ensure responsive design for all device types"""
    
    content.text = content_text
    
    # Slide 3: Key Features
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Features"
    content_text = """🔐 User Authentication
• Secure login/signup with validation
• Session management across browser sessions

🛍️ Product Catalog
• Browse products by categories
• Search functionality
• Detailed product information with image galleries

🛒 Shopping Cart
• Add, remove, and update product quantities
• Real-time cart updates
• Persistent cart data

📱 Responsive Design
• Mobile-first approach
• Modern UI with consistent design language
• Real-time validation with immediate feedback"""
    
    content.text = content_text
    
    # Slide 4: Technology Stack
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Stack"
    content_text = """Frontend Framework: React 18 with TypeScript
Build Tool: Vite for fast development and building
Styling: Tailwind CSS with Shadcn/ui components
State Management: React Context API
Routing: React Router DOM
Icons: Lucide React
Validation: Custom validation utilities

Benefits:
• Type safety with TypeScript
• Fast development with Vite
• Modern UI components with Shadcn/ui
• Efficient state management"""
    
    content.text = content_text
    
    # Slide 5: Functional Requirements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Functional Requirements"
    content_text = """👤 User Management
• User registration with email, name, and password
• User login/logout functionality
• Session persistence across browser sessions
• Authentication required for cart operations

📦 Product Management
• Browse all available products
• Filter products by category
• Search products by name
• View detailed product information

🛒 Shopping Cart
• Add products to cart (authenticated users)
• View cart contents and total
• Update product quantities
• Remove products and clear entire cart
• Cart persistence across sessions"""
    
    content.text = content_text
    
    # Slide 6: Non-Functional Requirements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Non-Functional Requirements"
    content_text = """⚡ Performance
• Application loads within 3 seconds
• Smooth animations and transitions (60fps)
• Efficient state management with minimal re-renders

🔒 Security
• Secure password requirements (8+ chars, complexity)
• Email format validation
• Input sanitization and validation
• Authentication required for sensitive operations

🎨 Usability
• Intuitive navigation and user flow
• Consistent design language throughout
• Accessible design with proper contrast
• Mobile-first responsive design

🛡️ Reliability
• Graceful error handling and recovery
• Data persistence across browser sessions
• Consistent behavior across different browsers"""
    
    content.text = content_text
    
    # Slide 7: Development Methodology
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Development Methodology"
    content_text = """🔄 Agile Development Approach
• Iterative development cycles
• Focus on delivering working features incrementally

📋 Development Process:
1. Planning Phase
   • Requirements analysis
   • Technology selection
   • Architecture planning
   • Timeline planning

2. Design Phase
   • UI/UX design with wireframes
   • Component design
   • State management design
   • API design

3. Development Phase
   • Sprint planning
   • Component development
   • Integration
   • Continuous testing

4. Testing Phase
   • Unit testing
   • Integration testing
   • User acceptance testing
   • Performance testing"""
    
    content.text = content_text
    
    # Slide 8: System Architecture
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Architecture"
    content_text = """🏗️ Component-Based Architecture

Presentation Layer: React components and UI
Business Logic Layer: Custom hooks and utilities
State Management Layer: React Context API
Data Layer: Static product data and localStorage

📦 Component Architecture:
• Core Components: App, Layout, Header, Footer
• Page Components: Index, Login, Signup, Products, ProductDetails, Cart
• UI Components: ProductCard, PasswordStrengthIndicator, Categories, FeaturedProducts

🔄 State Management:
• AuthContext: Manages user authentication state
• CartContext: Manages shopping cart state
• Local Storage: Persists user sessions and cart data

📊 Data Flow:
• User Actions → Component Events
• Context Providers → Application State
• Components → Re-render with new state
• UI → Updates reflect state changes"""
    
    content.text = content_text
    
    # Slide 9: Key Implementation Details
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Implementation Details"
    content_text = """🔐 Authentication System
• React Context for state management
• Secure login/signup with validation
• Session persistence using localStorage

🛒 Cart Management
• Efficient cart state updates
• Quantity management
• Cart persistence across sessions

✅ Validation System
• Custom validation utilities
• Real-time form validation
• Email format validation with regex
• Password strength requirements

📁 Code Organization
• src/components/: Reusable UI components
• src/pages/: Page-level components
• src/context/: React Context providers
• src/lib/: Utility functions and helpers
• src/hooks/: Custom React hooks
• src/data/: Static data and mock data"""
    
    content.text = content_text
    
    # Slide 10: Testing Strategy
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Testing Strategy"
    content_text = """🧪 Comprehensive Testing Approach

Unit Testing:
• Component functionality testing
• Utility function testing
• Custom hook testing

Integration Testing:
• Component interaction testing
• Context integration testing
• Routing functionality testing

User Acceptance Testing:
• Complete user journey testing
• Cross-browser compatibility
• Responsive functionality testing

🔍 Test Cases:
• Authentication: Valid/invalid login, registration, session persistence
• Cart Functionality: Add, update, remove items, cart persistence
• Product Browsing: Listing, filtering, search, details view
• Validation: Email, password, name validation with real-time feedback

⚡ Performance Testing:
• Load time within 3 seconds
• Smooth interactions and animations
• Efficient memory management
• Optimized JavaScript bundle"""
    
    content.text = content_text
    
    # Slide 11: Project Achievements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Achievements"
    content_text = """✅ Complete E-commerce Functionality
• Full product browsing, cart management, and user authentication
• Modern technology stack with React 18, TypeScript, Vite, and Tailwind CSS
• Responsive design with mobile-first approach

🎨 Technical Highlights
• Component-based architecture for modular and maintainable code
• Full TypeScript implementation for type safety
• Efficient state management using React Context API
• Custom validation system with comprehensive error handling
• Modern UI/UX with accessibility considerations

💼 Business Value
• Intuitive and engaging shopping experience
• Scalable architecture supporting future feature additions
• Clean code structure for easy maintenance
• Optimized performance for fast loading and smooth interactions
• Secure user data handling and authentication"""
    
    content.text = content_text
    
    # Slide 12: Lessons Learned
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Lessons Learned"
    content_text = """💡 Technical Lessons

State Management:
• Proper state management is crucial for complex applications
• React Context API provides excellent state management for medium-sized applications

Component Design:
• Reusable components significantly improve development efficiency
• Component composition and prop drilling alternatives are essential

TypeScript Benefits:
• TypeScript provides significant benefits in large projects
• Type safety reduces bugs and improves development experience

Performance Optimization:
• Performance considerations should be built into the development process
• Efficient re-renders and state updates are crucial for smooth UX

🔄 Development Process Lessons

Planning and Architecture:
• Proper planning and architecture design saves significant development time
• Component structure and data flow should be planned before implementation

Testing Strategy:
• Testing should be integrated into the development process
• Different types of testing are needed for different aspects of the application

User Experience:
• User experience should be prioritized from the beginning
• Small details like loading states and error handling significantly impact UX"""
    
    content.text = content_text
    
    # Slide 13: Future Improvements
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Improvements"
    content_text = """🚀 Technical Enhancements

Backend Integration:
• Implement real backend API integration
• Add proper database for product and user data
• Implement secure payment processing

Advanced Features:
• Advanced search and filtering capabilities
• Product review and rating system
• User wishlist functionality
• Complete order processing system
• Admin interface for product management

⚡ Performance Optimizations

Code Splitting:
• Implement lazy loading for better performance
• Add proper caching strategies
• Implement content delivery network
• Convert to Progressive Web Application

📱 Additional Features:
• User Reviews and Ratings
• Wishlist Functionality
• Order Management System
• Admin Panel for Product Management
• Advanced Analytics and Reporting"""
    
    content.text = content_text
    
    # Slide 14: Conclusion
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    content_text = """🎯 Project Success

The Shop Hub e-commerce application successfully demonstrates:
• Modern web development practices
• Solid foundation for production-ready e-commerce platform
• Importance of proper planning and architecture design
• User experience considerations in web application development

📚 Key Takeaways

• Lessons learned will be valuable for future development projects
• Insights into building scalable, maintainable, and user-friendly web applications
• Understanding of modern React development patterns
• Experience with TypeScript and modern build tools

🔗 Project Links

GitHub Repository: https://github.com/vedant208purohit/shop-hub-208.git
Live Demo: [Available upon request]
Demo Credentials: demo@example.com / password

Thank You!"""
    
    content.text = content_text
    
    # Save the presentation
    prs.save('Shop_Hub_Presentation.pptx')
    print("PowerPoint presentation created successfully: Shop_Hub_Presentation.pptx")

if __name__ == "__main__":
    create_shop_hub_presentation() 